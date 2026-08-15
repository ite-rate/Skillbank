#!/usr/bin/env python3
"""Structural QA for generated .pptx decks.

Usage: python3 qa_deck.py <deck.pptx> [expected_pages]

Checks (read-only, no rendering needed):
  - File opens cleanly with python-pptx (catches XML corruption)
  - Slide count vs expected
  - Every shape's bounding box inside the canvas (catches out-of-bounds layout bugs)
  - Per-slide shape/text counts + preview, so empty slides are obvious

Pairs with content QA: `markitdown deck.pptx | grep -i "object Object"`
which catches rich-text corruption that this script cannot see.

Requires: pip install python-pptx
"""
import sys
from pptx import Presentation
from pptx.util import Emu

path = sys.argv[1] if len(sys.argv) > 1 else "deck.pptx"
expected = int(sys.argv[2]) if len(sys.argv) > 2 else None

p = Presentation(path)
W, H = p.slide_width, p.slide_height
w_in, h_in = Emu(W).inches, Emu(H).inches
print(f"canvas: {w_in:.2f} x {h_in:.2f} in, slides: {len(p.slides)}")
if expected and len(p.slides) != expected:
    print(f"!! EXPECTED {expected} slides, got {len(p.slides)}")

problems = []
for i, slide in enumerate(p.slides, 1):
    n_shapes, chars = 0, 0
    preview = ""
    for sh in slide.shapes:
        n_shapes += 1
        try:
            x, y, w, h = sh.left, sh.top, sh.width, sh.height
            if None not in (x, y, w, h):
                xi, yi, wi, hi = Emu(x).inches, Emu(y).inches, Emu(w).inches, Emu(h).inches
                if xi < -0.02 or yi < -0.02 or xi + wi > w_in + 0.02 or yi + hi > h_in + 0.02:
                    problems.append(
                        f"P{i} out-of-bounds {sh.shape_type} @({xi:.2f},{yi:.2f}) {wi:.2f}x{hi:.2f}")
        except Exception:
            pass
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                chars += len(t)
                if not preview:
                    preview = t[:30]
    print(f"P{i:>2}: {n_shapes:>3} shapes, {chars:>4} chars, {preview}")

print("\n== issues ==")
print("\n".join(problems) if problems else "none")
